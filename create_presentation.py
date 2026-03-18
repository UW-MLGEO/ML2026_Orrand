"""
Generate a PowerPoint presentation for the pCO2 ML Prediction project.
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── paths ──────────────────────────────────────────────────────────────
BASE = os.path.dirname(os.path.abspath(__file__))
PLOTS = os.path.join(BASE, "plots")
ML = os.path.join(PLOTS, "ml_results")
EDA = os.path.join(PLOTS, "training_data_eda")
ANALYSIS = os.path.join(PLOTS, "analysis")
EXPLORATION = os.path.join(PLOTS, "exploration")
OUTPUT = os.path.join(BASE, "pCO2_ML_Presentation.pptx")

# ── colours ────────────────────────────────────────────────────────────
DARK_BLUE  = RGBColor(0x00, 0x2B, 0x5C)   # navy header / title
MID_BLUE   = RGBColor(0x00, 0x6D, 0xAA)   # accent
LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)   # slide bg
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
BLACK      = RGBColor(0x00, 0x00, 0x00)
OCEAN_TEAL = RGBColor(0x00, 0x89, 0x9B)

SLIDE_W = Inches(13.333)  # widescreen 16:9
SLIDE_H = Inches(7.5)

# ── helpers ────────────────────────────────────────────────────────────
def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_size=18,
                bold=False, color=BLACK, alignment=PP_ALIGN.LEFT,
                font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return tf


def add_bullet_slide_content(tf, bullets, font_size=16, color=BLACK,
                              font_name="Calibri", bold=False, spacing=Pt(6)):
    """Add bullet points to an existing text frame."""
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = font_name
        p.font.bold = bold
        p.level = 0
        p.space_after = spacing
    return tf


def add_image_safe(slide, img_path, left, top, width=None, height=None):
    """Add image only if file exists, else add a placeholder box."""
    if os.path.isfile(img_path):
        if width and height:
            slide.shapes.add_picture(img_path, left, top, width, height)
        elif width:
            slide.shapes.add_picture(img_path, left, top, width=width)
        elif height:
            slide.shapes.add_picture(img_path, left, top, height=height)
        else:
            slide.shapes.add_picture(img_path, left, top)
    else:
        # draw a light-gray rectangle placeholder
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, left, top,
            width or Inches(5), height or Inches(3.5))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"[Image not found]\n{os.path.basename(img_path)}"
        p.font.size = Pt(10)
        p.font.color.rgb = RGBColor(0x99, 0x99, 0x99)


def make_title_bar(slide, title_text, subtitle_text=None):
    """Dark-blue banner across the top with white title text."""
    # banner shape
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, Inches(1.15))
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BLUE
    shape.line.fill.background()
    # title
    add_textbox(slide, Inches(0.6), Inches(0.15), Inches(12), Inches(0.6),
                title_text, font_size=28, bold=True, color=WHITE)
    if subtitle_text:
        add_textbox(slide, Inches(0.6), Inches(0.7), Inches(12), Inches(0.4),
                    subtitle_text, font_size=14, color=RGBColor(0xBB, 0xCC, 0xDD))


# ── build presentation ─────────────────────────────────────────────────
prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
blank_layout = prs.slide_layouts[6]  # blank

# ====================================================================
# SLIDE 1 – Title
# ====================================================================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, DARK_BLUE)

add_textbox(slide, Inches(1), Inches(1.8), Inches(11), Inches(1.2),
            "Predicting Ocean pCO\u2082\nUsing Machine Learning",
            font_size=40, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(1), Inches(3.5), Inches(11), Inches(0.5),
            "Bridging Sparse Buoy Measurements with Satellite Observations",
            font_size=20, color=RGBColor(0xAA, 0xCC, 0xEE), alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(1), Inches(5.0), Inches(11), Inches(0.8),
            "MLGEO 2026  |  University of Washington\nMarch 2026",
            font_size=16, color=RGBColor(0x88, 0xAA, 0xCC), alignment=PP_ALIGN.CENTER)

# thin accent line
shape = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(4), Inches(4.2), Inches(5.333), Inches(0.04))
shape.fill.solid()
shape.fill.fore_color.rgb = OCEAN_TEAL
shape.line.fill.background()


# ====================================================================
# SLIDE 2 – Motivation
# ====================================================================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, WHITE)
make_title_bar(slide, "Why Predict Ocean pCO\u2082?")

bullets = [
    "The ocean absorbs ~25% of anthropogenic CO\u2082 \u2014 understanding this sink is critical for climate projections.",
    "pCO\u2082 (partial pressure of CO\u2082 in seawater) quantifies air-sea CO\u2082 exchange.",
    "Direct measurements are sparse: limited to ~7 NOAA moored buoys in U.S. waters.",
    "Satellites observe SST and chlorophyll-a globally and daily \u2014 but not pCO\u2082.",
    "Goal: Use ML to predict pCO\u2082 from satellite-derived features, enabling ocean-wide estimates.",
]
tf = add_textbox(slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(5.0),
                 "", font_size=18)
add_bullet_slide_content(tf, bullets, font_size=19, spacing=Pt(12))


# ====================================================================
# SLIDE 3 – The Challenge  (with records-per-location figure)
# ====================================================================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, WHITE)
make_title_bar(slide, "The Challenge: Sparse Observations")

add_textbox(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(4.5),
            "Only 7 buoy sites provide continuous pCO\u2082 measurements.\n\n"
            "Each site has different lengths of record and gaps in coverage.\n\n"
            "We need models that generalize across diverse ocean regimes \u2014 "
            "from the Bering Sea to the South Pacific.",
            font_size=17)

add_image_safe(slide,
               os.path.join(PLOTS, "presentation_records_per_location.png"),
               Inches(6.8), Inches(1.4), width=Inches(5.8))


# ====================================================================
# SLIDE 4 – Satellite Data Source
# ====================================================================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, WHITE)
make_title_bar(slide, "Data: Satellite Observations", "MUR SST & MODIS Chlorophyll-a")

add_textbox(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(5.0),
            "Sea Surface Temperature (SST)\n"
            "\u2022  MUR L4 product (1 km, daily)\n"
            "\u2022  7 statistical features per buoy-day\n\n"
            "Chlorophyll-a (chl-a)\n"
            "\u2022  MODIS Aqua (4 km, 8-day composites)\n"
            "\u2022  8 features incl. temporal offset\n\n"
            "Both extracted for a pixel window\naround each buoy site via ERDDAP.",
            font_size=16)

add_image_safe(slide,
               os.path.join(PLOTS, "01_satellite_sst_overview.png"),
               Inches(6.8), Inches(1.4), width=Inches(5.8))


# ====================================================================
# SLIDE 5 – Buoy Data
# ====================================================================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, WHITE)
make_title_bar(slide, "Data: NOAA Buoy Observations")

add_textbox(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(5.0),
            "7 moored ocean buoy sites across U.S. waters:\n\n"
            "\u2022  In-situ SST (for satellite validation)\n"
            "\u2022  In-situ pCO\u2082 (our prediction target)\n"
            "\u2022  Latitude & longitude\n\n"
            "Data aggregated to daily means.\n"
            "Quality-controlled for sensor drift & missing values.",
            font_size=17)

add_image_safe(slide,
               os.path.join(PLOTS, "02_buoy_data_overview.png"),
               Inches(6.8), Inches(1.4), width=Inches(5.8))


# ====================================================================
# SLIDE 6 – Data Availability Timeline
# ====================================================================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, WHITE)
make_title_bar(slide, "Data Availability & Coverage")

add_image_safe(slide,
               os.path.join(PLOTS, "05_pco2_data_availability_timeline.png"),
               Inches(0.8), Inches(1.5), width=Inches(11.7))

add_textbox(slide, Inches(0.8), Inches(6.0), Inches(11), Inches(0.8),
            "479 buoy-days of matched satellite + buoy data across 7 sites.",
            font_size=15, color=RGBColor(0x55, 0x55, 0x55))


# ====================================================================
# SLIDE 7 – Dataset Dashboard
# ====================================================================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, WHITE)
make_title_bar(slide, "Constructed ML Dataset")

add_image_safe(slide,
               os.path.join(PLOTS, "03_ml_dataset_dashboard.png"),
               Inches(0.5), Inches(1.4), width=Inches(12.3))

add_textbox(slide, Inches(0.8), Inches(6.2), Inches(11), Inches(0.8),
            "18 input features  |  1 target (daily mean pCO\u2082)  |  479 samples from 7 sites",
            font_size=15, color=RGBColor(0x55, 0x55, 0x55), alignment=PP_ALIGN.CENTER)


# ====================================================================
# SLIDE 8 – EDA: pCO2 variability by site
# ====================================================================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, WHITE)
make_title_bar(slide, "Exploratory Data Analysis", "pCO\u2082 variability across buoy sites")

add_textbox(slide, Inches(0.8), Inches(1.5), Inches(5.0), Inches(3.0),
            "Wide range of pCO\u2082 across sites:\n\n"
            "\u2022 Some sites cluster near atmospheric (~400 \u00b5atm)\n"
            "\u2022 Others show large seasonal swings\n"
            "\u2022 This variability is what the model must capture",
            font_size=16)

add_image_safe(slide,
               os.path.join(EDA, "03_pco2_by_site_boxplot.png"),
               Inches(6.2), Inches(1.4), width=Inches(6.5))


# ====================================================================
# SLIDE 9 – EDA: Feature-target relationships
# ====================================================================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, WHITE)
make_title_bar(slide, "Feature-Target Relationships", "SST and Chl-a vs. pCO\u2082")

add_image_safe(slide,
               os.path.join(EDA, "04_sst_chla_vs_pco2.png"),
               Inches(0.5), Inches(1.4), width=Inches(12.3))

add_textbox(slide, Inches(0.8), Inches(6.2), Inches(11), Inches(0.8),
            "Non-linear relationships between features and pCO\u2082 motivate tree-based models.",
            font_size=15, color=RGBColor(0x55, 0x55, 0x55), alignment=PP_ALIGN.CENTER)


# ====================================================================
# SLIDE 10 – Methods Overview
# ====================================================================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, WHITE)
make_title_bar(slide, "Machine Learning Approach")

bullets_left = [
    "Train / Test Split",
    "\u2022  80 / 20 stratified by site",
    "\u2022  Features z-score standardized",
    "",
    "Two Models of Increasing Complexity:",
    "1. Linear Regression (baseline)",
    "2. Random Forest (200 trees)",
]
tf = add_textbox(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(5.0),
                 "", font_size=17)
add_bullet_slide_content(tf, bullets_left, font_size=17, spacing=Pt(8))

# Evaluation metrics on right side
bullets_right = [
    "Evaluation Metrics:",
    "",
    "RMSE \u2013 root mean squared error (\u00b5atm)",
    "MAE  \u2013 mean absolute error (\u00b5atm)",
    "R\u00b2   \u2013 variance explained (1.0 = perfect)",
]
tf2 = add_textbox(slide, Inches(7.0), Inches(1.5), Inches(5.5), Inches(5.0),
                  "", font_size=17)
add_bullet_slide_content(tf2, bullets_right, font_size=17, spacing=Pt(8))


# ====================================================================
# SLIDE 11 – Model 1: Linear Regression
# ====================================================================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, WHITE)
make_title_bar(slide, "Model 1: Linear Regression", "The Simple Baseline")

bullets = [
    "Fits a single linear equation:",
    "    pCO\u2082 = a\u00d7SST + b\u00d7chl-a + c\u00d7lat + \u2026",
    "",
    "Strengths:",
    "\u2022  Fast, interpretable, good starting point",
    "",
    "Weaknesses:",
    "\u2022  Assumes linear feature\u2013target relationships",
    "\u2022  Ocean carbon chemistry is highly non-linear",
]
tf = add_textbox(slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(5.0),
                 "", font_size=18)
add_bullet_slide_content(tf, bullets, font_size=18, spacing=Pt(8))


# ====================================================================
# SLIDE 12 – Model 2: Random Forest
# ====================================================================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, WHITE)
make_title_bar(slide, "Model 2: Random Forest", "Ensemble of 200 Decision Trees")

bullets = [
    "Builds 200 decision trees, each on a random subset of data & features.",
    "Final prediction = average of all trees.",
    "",
    "Strengths:",
    "\u2022  Captures non-linear relationships naturally",
    "\u2022  Robust, hard to overfit",
    "",
    "Weaknesses:",
    "\u2022  Slower on very large datasets",
    "\u2022  Does not extrapolate beyond training data range",
]
tf = add_textbox(slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(5.0),
                 "", font_size=18)
add_bullet_slide_content(tf, bullets, font_size=18, spacing=Pt(8))


# ====================================================================
# SLIDE 13 – Results: Predicted vs Actual
# ====================================================================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, WHITE)
make_title_bar(slide, "Results: Predicted vs. Actual pCO\u2082")

add_image_safe(slide,
               os.path.join(ML, "01_predicted_vs_actual.png"),
               Inches(0.5), Inches(1.4), width=Inches(12.3))

add_textbox(slide, Inches(0.8), Inches(6.2), Inches(11), Inches(0.8),
            "Tree-based models cluster tightly around the 1:1 line; Linear Regression shows much more scatter.",
            font_size=15, color=RGBColor(0x55, 0x55, 0x55), alignment=PP_ALIGN.CENTER)


# ====================================================================
# SLIDE 15 – Results: Residuals
# ====================================================================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, WHITE)
make_title_bar(slide, "Results: Residual Analysis")

add_image_safe(slide,
               os.path.join(ML, "02_residuals.png"),
               Inches(0.5), Inches(1.4), width=Inches(12.3))

add_textbox(slide, Inches(0.8), Inches(6.2), Inches(11), Inches(0.8),
            "Residuals centered around zero with no systematic bias. Most errors are small.",
            font_size=15, color=RGBColor(0x55, 0x55, 0x55), alignment=PP_ALIGN.CENTER)


# ====================================================================
# SLIDE 16 – Results: Feature Importance
# ====================================================================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, WHITE)
make_title_bar(slide, "Results: Feature Importance")

add_textbox(slide, Inches(0.8), Inches(1.5), Inches(5.0), Inches(4.5),
            "SST features dominate \u2014 consistent with\n"
            "known ocean chemistry (CO\u2082 solubility\n"
            "decreases as water warms).\n\n"
            "Location (lat/lon) ranks high, reflecting\n"
            "regional ocean carbon regimes.\n\n"
            "Chlorophyll-a provides additional\n"
            "information on biological carbon uptake.",
            font_size=16)

add_image_safe(slide,
               os.path.join(ML, "03_feature_importance.png"),
               Inches(6.2), Inches(1.4), width=Inches(6.5))


# ====================================================================
# SLIDE 17 – Results: Performance by Site
# ====================================================================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, WHITE)
make_title_bar(slide, "Results: Performance by Buoy Site")

add_textbox(slide, Inches(0.8), Inches(1.5), Inches(5.0), Inches(4.5),
            "Model skill varies across sites:\n\n"
            "\u2022 Open-ocean sites: more predictable\n"
            "\u2022 Coastal/estuary sites: higher error\n\n"
            "Sites with more scatter may need\n"
            "additional features or more training data.",
            font_size=16)

add_image_safe(slide,
               os.path.join(ML, "04_predictions_by_site.png"),
               Inches(6.2), Inches(1.4), width=Inches(6.5))


# ====================================================================
# SLIDE 18 – Summary 6-Panel Figure
# ====================================================================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, WHITE)
make_title_bar(slide, "Project Summary Overview")

add_image_safe(slide,
               os.path.join(PLOTS, "summary_figure_6panel.png"),
               Inches(0.5), Inches(1.3), width=Inches(12.3))


# ====================================================================
# SLIDE 19 – Key Takeaways
# ====================================================================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, WHITE)
make_title_bar(slide, "Key Takeaways")

bullets = [
    "1.  Tree-based models significantly outperform Linear Regression \u2014 ocean pCO\u2082 relationships are non-linear.",
    "",
    "2.  Sea surface temperature is the most important predictor, consistent with CO\u2082 solubility physics.",
    "",
    "3.  Some buoy sites are harder to predict, reflecting real physical differences in carbon dynamics.",
    "",
    "4.  ML can bridge the gap between sparse buoy measurements and continuous satellite coverage,\n"
    "     enabling potential ocean-wide pCO\u2082 estimates.",
]
tf = add_textbox(slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(5.5),
                 "", font_size=19)
add_bullet_slide_content(tf, bullets, font_size=19, spacing=Pt(8))


# ====================================================================
# SLIDE 20 – Next Steps
# ====================================================================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, WHITE)
make_title_bar(slide, "Next Steps")

bullets = [
    "\u2022  Hyperparameter tuning via cross-validation / Bayesian optimization",
    "\u2022  Add more satellite features (salinity, wind speed, mixed layer depth)",
    "\u2022  Incorporate temporal features (day-of-year, seasonal indicators)",
    "\u2022  K-fold cross-validation for robust performance estimates",
    "\u2022  Leave-one-site-out experiment to test spatial generalization",
    "\u2022  Explore neural network architectures",
    "\u2022  Deploy model on full satellite grids to produce spatial pCO\u2082 maps",
]
tf = add_textbox(slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(5.5),
                 "", font_size=19)
add_bullet_slide_content(tf, bullets, font_size=19, spacing=Pt(14))


# ====================================================================
# SLIDE 21 – Thank You / Questions
# ====================================================================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, DARK_BLUE)

add_textbox(slide, Inches(1), Inches(2.5), Inches(11), Inches(1.0),
            "Thank You!", font_size=44, bold=True, color=WHITE,
            alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(1), Inches(4.0), Inches(11), Inches(0.6),
            "Questions?", font_size=28, color=RGBColor(0xAA, 0xCC, 0xEE),
            alignment=PP_ALIGN.CENTER)

shape = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(4), Inches(3.7), Inches(5.333), Inches(0.04))
shape.fill.solid()
shape.fill.fore_color.rgb = OCEAN_TEAL
shape.line.fill.background()


# ── save ───────────────────────────────────────────────────────────────
prs.save(OUTPUT)
print(f"\nPresentation saved to:\n  {OUTPUT}")
print(f"  ({len(prs.slides)} slides)")
